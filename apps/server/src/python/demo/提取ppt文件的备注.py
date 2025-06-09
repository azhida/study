from pptx import Presentation

# 提取 PPT 文件备注内容
def extract_notes(ppt_path):
  prs = Presentation(ppt_path)
  notes = []
  for slide in prs.slides:
    notes.append(slide.notes_slide.notes_text_frame.text)
  return notes
    
motes = extract_notes('./demo.pptx')
print(motes)

